"""Raw text extraction for the supported formats (PRD F4, F5).

PDF (text layer), Markdown/text variants, and HTML (fetched pages or .html
files). OCR for scanned PDFs is out of scope (PRD §2).
"""
from __future__ import annotations

import io
import re
from pathlib import Path

import requests
from bs4 import BeautifulSoup

TEXT_EXTENSIONS = {
    ".md", ".txt", ".markdown", ".rst", ".csv", ".json", ".yaml", ".yml", ".log",
}
HTML_EXTENSIONS = {".html", ".htm"}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
PPTX_EXTENSIONS = {".pptx"}
XLSX_EXTENSIONS = {".xlsx", ".xlsm"}


def load_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            parts.append(text)
    return "\n\n".join(parts)


def load_html(data: bytes, url: str | None = None) -> str:
    soup = BeautifulSoup(data, "html.parser")
    for tag in soup(["script", "style", "noscript", "svg"]):
        tag.decompose()
    title = soup.title.get_text(strip=True) if soup.title else ""
    text = soup.get_text(separator="\n")
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if title:
        text = f"{title}\n\n{text}"
    return text


def load_docx(data: bytes) -> str:
    """Extract text from a .docx (Word) file, preserving paragraph breaks."""
    from docx import Document

    doc = Document(io.BytesIO(data))
    parts = []
    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    for table in doc.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            line = " | ".join(c for c in cells if c)
            if line:
                parts.append(line)
    return "\n\n".join(parts)


def load_pptx(data: bytes) -> str:
    """Extract text from a .pptx (PowerPoint) file, slide by slide."""
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        parts = [f"# Slide {i}"]
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    line = " | ".join(c for c in cells if c)
                    if line:
                        parts.append(line)
        if len(parts) > 1:
            slides.append("\n".join(parts))
    return "\n\n".join(slides)


def load_xlsx(data: bytes) -> str:
    """Dump every sheet of an .xlsx workbook to delimited text."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if v is None else str(v) for v in row]
            line = " | ".join(cells).strip()
            if line:
                rows.append(line)
        if rows:
            sheets.append(f"# Sheet: {ws.title}\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def load_bytes(filename: str, data: bytes, url: str | None = None) -> str:
    ext = Path(filename).suffix.lower()
    if ext in PDF_EXTENSIONS:
        return load_pdf(data)
    if ext in HTML_EXTENSIONS or url is not None:
        return load_html(data, url)
    if ext in DOCX_EXTENSIONS:
        return load_docx(data)
    if ext in PPTX_EXTENSIONS:
        return load_pptx(data)
    if ext in XLSX_EXTENSIONS:
        return load_xlsx(data)
    if ext in TEXT_EXTENSIONS or ext == "":
        return data.decode("utf-8", errors="replace")
    raise ValueError(f"Unsupported file type: {ext or filename}")


def fetch_url(url: str) -> tuple[str, bytes]:
    """Fetch a web page, following redirects. Returns (final_url, body)."""
    resp = requests.get(
        url,
        timeout=30,
        headers={"User-Agent": "ragchat/1.0 (+class project)"},
        allow_redirects=True,
    )
    resp.raise_for_status()
    return resp.url, resp.content


def page_title(url: str, body: bytes) -> str:
    soup = BeautifulSoup(body, "html.parser")
    if soup.title and soup.title.get_text(strip=True):
        return soup.title.get_text(strip=True)[:200]
    return url[:200]
